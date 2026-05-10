"""Render a BlocksDocument to a .pptx via python-pptx.

Slide mapping:
- Each Lesson opens a "title slide" (using lesson.title).
- Each subsequent `title` block within the lesson starts a new content slide.
- `text`, `table`, and `image` blocks are appended to the current content slide.
"""
from __future__ import annotations

import io
import json
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from officeplane.content_agent.renderers.blocks import (
    Block,
    BlocksDocument,
    Lesson,
)

log = logging.getLogger("officeplane.content_agent.renderers.pptx_blocks")


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CONTENT_TOP = Inches(1.5)
CONTENT_LEFT = Inches(0.6)
CONTENT_WIDTH = Inches(12.1)


def render_pptx(doc: BlocksDocument, *, image_root: Path | None = None) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover slide
    _cover_slide(prs, doc.title)

    for module in doc.modules:
        if module.title:
            _section_slide(prs, module.title)
        for lesson in module.lessons:
            _render_lesson(prs, lesson, image_root=image_root)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _cover_slide(prs, title: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title or "Untitled"


def _section_slide(prs, section_title: str) -> None:
    # Use layout index 2 if available (section header), otherwise fall back to 0
    try:
        layout = prs.slide_layouts[2]
    except IndexError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = section_title


def _render_lesson(prs, lesson: Lesson, *, image_root: Path | None) -> None:
    """One title slide for the lesson + one content slide per `title` block boundary."""
    if lesson.title:
        _content_slide(prs, lesson.title, body_blocks=[])

    current_title: str | None = lesson.title or None
    pending: list[Block] = []
    for block in sorted(lesson.blocks, key=lambda b: b.order):
        if block.type == "title":
            # Flush previous slide
            if pending or current_title != lesson.title:
                _content_slide(prs, current_title or "", body_blocks=pending, image_root=image_root)
            current_title = block.content or ""
            pending = []
        else:
            pending.append(block)
    if pending:
        _content_slide(prs, current_title or "", body_blocks=pending, image_root=image_root)


def _content_slide(prs, title: str, body_blocks: list[Block], *, image_root: Path | None = None) -> None:
    # Use layout index 5 if available (blank/title-only), fall back to 1
    try:
        layout = prs.slide_layouts[5]
    except IndexError:
        layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title or " "

    if not body_blocks:
        return

    # Compose body content
    has_table = any(b.type == "table" for b in body_blocks)
    has_image = any(b.type == "image" for b in body_blocks)

    if has_table:
        spec = next((json.loads(b.content) for b in body_blocks if b.type == "table" and b.content), None)
        if spec:
            _add_table(slide, spec)

    if has_image and image_root is not None:
        for b in body_blocks:
            if b.type == "image" and b.object_key:
                path = image_root / b.object_key
                if path.exists():
                    slide.shapes.add_picture(str(path), Inches(0.8), Inches(2.0), height=Inches(4.5))

    # Text blocks → bullets in a text box
    text_blocks = [b for b in body_blocks if b.type == "text" and b.content]
    if text_blocks:
        textbox = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(5.5))
        tf = textbox.text_frame
        tf.word_wrap = True
        for i, b in enumerate(text_blocks):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b.content or ""
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)


def _add_table(slide, spec: dict) -> None:
    headers = list(spec.get("headers") or [])
    rows = [list(r) for r in (spec.get("rows") or [])]
    if not headers and not rows:
        return
    cols = max(len(headers), max((len(r) for r in rows), default=0))
    if cols == 0:
        return
    total_rows = 1 + len(rows) if headers else len(rows)
    if total_rows == 0:
        return
    table_shape = slide.shapes.add_table(
        total_rows, cols, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, Inches(4.0)
    )
    table = table_shape.table
    if headers:
        for c, h in enumerate(headers):
            table.cell(0, c).text = str(h)
        offset = 1
    else:
        offset = 0
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            table.cell(r_idx + offset, c_idx).text = str(val)
