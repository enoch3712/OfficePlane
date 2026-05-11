"""PPTX renderer that walks the agnostic Document tree and produces .pptx bytes.

Uses python-pptx to emit a well-structured PowerPoint presentation.  The
renderer is intentionally stateless — each call to :func:`render_pptx` opens a
fresh :class:`~pptx.Presentation`, collects a flat list of *slide specs*, then
emits slides in order.

Slide layout strategy
---------------------
1. **Title slide** (always first if ``doc.meta.title`` is non-empty):
   layout 0 (``Title Slide``).

2. **Section divider slide** for each top-level (level=1) Section that has
   child Sections beneath it: layout 2 (``Section Header``) if available,
   else layout 1.

3. **Content slide** for each leaf-bearing section (a Section whose
   ``children`` contain at least one non-Section block): layout 1
   (``Title and Content``).

4. **Big-primitive slides** (Table / Figure / Code) each get their own slide.

The renderer batches consecutive small blocks (heading / paragraph / list /
callout / quote / divider) into the same content slide body, flushing on each
big block or nested Section.
"""

from __future__ import annotations

import logging
from io import BytesIO
from os.path import isfile
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # noqa: F401 — imported for future use
from pptx.util import Inches, Pt

from officeplane.content_agent.renderers.document import (
    Block,
    Callout,
    Code,
    Divider,
    Document,
    Figure,
    Heading,
    List,
    Paragraph,
    Quote,
    Section,
    Table,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Slide geometry
# ---------------------------------------------------------------------------

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)

# Layout indices by purpose (fallbacks handled at runtime)
_LAYOUT_TITLE_SLIDE = 0
_LAYOUT_TITLE_CONTENT = 1
_LAYOUT_SECTION_HEADER = 2

# ---------------------------------------------------------------------------
# Slide spec dataclasses (plain dicts are fine given the internal scope)
# ---------------------------------------------------------------------------

_TITLE_SLIDE = "title_slide"
_SECTION_DIVIDER = "section_divider"
_CONTENT_SLIDE = "content_slide"
_TABLE_SLIDE = "table_slide"
_FIGURE_SLIDE = "figure_slide"
_CODE_SLIDE = "code_slide"

SlideSpec = dict[str, Any]


# ---------------------------------------------------------------------------
# Spec collection — walk the tree once, build a flat list of SlideSpec dicts
# ---------------------------------------------------------------------------


def _is_big_block(node: Section | Block) -> bool:
    """True for block types that require their own slide."""
    return isinstance(node, (Table, Figure, Code))


def _collect_specs(doc: Document) -> list[SlideSpec]:
    specs: list[SlideSpec] = []

    # Title slide
    if doc.meta.title:
        # Find first top-level section to use as subtitle
        subtitle = ""
        for child in doc.children:
            if isinstance(child, Section):
                subtitle = child.heading
                break
        specs.append(
            {"kind": _TITLE_SLIDE, "title": doc.meta.title, "subtitle": subtitle}
        )

    # Walk top-level children
    for child in doc.children:
        if isinstance(child, Section):
            _collect_section_specs(child, specs)
        else:
            # Top-level loose blocks go into a generic content slide
            _flush_block_to_specs(child, "Document", specs)

    return specs


def _has_sub_sections(section: Section) -> bool:
    """True if any direct child is itself a Section."""
    return any(isinstance(c, Section) for c in section.children)


def _collect_section_specs(section: Section, specs: list[SlideSpec]) -> None:
    """Recursively collect slide specs for a Section and its descendants."""
    # Section divider for non-leaf parent sections
    if section.level == 1 and _has_sub_sections(section):
        specs.append({"kind": _SECTION_DIVIDER, "heading": section.heading})
        # Recurse into child sections; loose blocks at this level are batched
        # into a content slide
        _process_mixed_children(section, specs)
        return

    # Leaf-bearing section (has at least one non-Section child) — content slide
    block_children = [c for c in section.children if not isinstance(c, Section)]
    if block_children:
        _process_mixed_children(section, specs)
    else:
        # All children are sub-sections — emit divider + recurse
        if section.children:
            specs.append({"kind": _SECTION_DIVIDER, "heading": section.heading})
            for child in section.children:
                if isinstance(child, Section):
                    _collect_section_specs(child, specs)


def _process_mixed_children(section: Section, specs: list[SlideSpec]) -> None:
    """Walk a section's children, batching small blocks and flushing on big ones."""
    pending_blocks: list[Block] = []  # accumulated small blocks

    def flush(heading: str = "") -> None:
        if pending_blocks:
            specs.append(
                {
                    "kind": _CONTENT_SLIDE,
                    "heading": heading,
                    "blocks": list(pending_blocks),
                }
            )
            pending_blocks.clear()

    current_heading = section.heading

    for child in section.children:
        if isinstance(child, Section):
            flush(current_heading)
            _collect_section_specs(child, specs)
        elif _is_big_block(child):
            flush(current_heading)
            _flush_block_to_specs(child, current_heading, specs)
        else:
            pending_blocks.append(child)  # type: ignore[arg-type]

    flush(current_heading)


def _flush_block_to_specs(
    block: Section | Block, heading: str, specs: list[SlideSpec]
) -> None:
    """Emit a standalone slide spec for a big block."""
    if isinstance(block, Table):
        specs.append({"kind": _TABLE_SLIDE, "heading": heading, "table": block})
    elif isinstance(block, Figure):
        specs.append({"kind": _FIGURE_SLIDE, "heading": heading, "figure": block})
    elif isinstance(block, Code):
        specs.append({"kind": _CODE_SLIDE, "heading": heading, "code": block})
    else:
        # Small block at top-level — wrap in a minimal content slide
        specs.append(
            {"kind": _CONTENT_SLIDE, "heading": heading, "blocks": [block]}
        )


# ---------------------------------------------------------------------------
# Slide emission — translate specs into actual pptx slides
# ---------------------------------------------------------------------------


def _safe_layout(prs: Presentation, idx: int) -> Any:
    """Return a slide layout by index, falling back to 0 if out-of-range."""
    layouts = prs.slide_layouts
    if idx < len(layouts):
        return layouts[idx]
    return layouts[0]


def _add_title_slide(prs: Presentation, spec: SlideSpec) -> None:
    layout = _safe_layout(prs, _LAYOUT_TITLE_SLIDE)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec["title"]
    # Subtitle placeholder is index 1 in Title Slide layout
    if spec.get("subtitle") and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = spec["subtitle"]
        except Exception:
            pass


def _add_section_divider(prs: Presentation, spec: SlideSpec) -> None:
    # Use layout 2 if available, else fall back to 1
    n = len(prs.slide_layouts)
    layout_idx = _LAYOUT_SECTION_HEADER if n > _LAYOUT_SECTION_HEADER else _LAYOUT_TITLE_CONTENT
    layout = _safe_layout(prs, layout_idx)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec["heading"]


def _add_content_slide(prs: Presentation, spec: SlideSpec) -> None:
    layout = _safe_layout(prs, _LAYOUT_TITLE_CONTENT)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.get("heading", "")

    # Find body placeholder (index 1)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break

    if body_ph is None:
        return

    tf = body_ph.text_frame
    tf.word_wrap = True

    blocks: list[Block] = spec.get("blocks", [])
    first = True
    for block in blocks:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        _populate_paragraph(para, block)


def _populate_paragraph(para: Any, block: Block) -> None:
    """Set text/style on a pptx paragraph based on block type."""
    if isinstance(block, Heading):
        para.text = block.text
        para.font.bold = True
        para.level = max(0, block.level - 2)
    elif isinstance(block, Paragraph):
        para.text = block.text
    elif isinstance(block, List):
        # First item goes into the provided paragraph, rest get new ones
        items = block.items
        if items:
            para.text = items[0].text
            para.level = 1
            tf = para._p.getparent().getparent()  # access to text frame — handled below
        # Note: list items beyond the first cannot be added here without the
        # text frame reference. We handle multi-item lists via a wrapper approach.
        # This is handled in _add_content_slide by passing each item separately.
    elif isinstance(block, Callout):
        para.text = f"[{block.variant.upper()}] {block.text}"
    elif isinstance(block, Quote):
        para.text = block.text
        para.font.italic = True
    elif isinstance(block, Divider):
        para.text = ""


def _populate_body(tf: Any, blocks: list[Block]) -> None:
    """Populate a text frame with multiple blocks, handling List multi-items."""
    first = True
    for block in blocks:
        if isinstance(block, List):
            for item in block.items:
                if first:
                    para = tf.paragraphs[0]
                    first = False
                else:
                    para = tf.add_paragraph()
                para.text = item.text
                para.level = 1
        elif isinstance(block, Heading):
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.text = block.text
            para.font.bold = True
            para.level = max(0, block.level - 2)
        elif isinstance(block, Paragraph):
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.text = block.text
        elif isinstance(block, Callout):
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.text = f"[{block.variant.upper()}] {block.text}"
        elif isinstance(block, Quote):
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.text = block.text
            para.font.italic = True
        elif isinstance(block, Divider):
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.text = ""
        else:
            logger.warning("pptx_render: unknown block type %s — skipped", type(block).__name__)


def _add_content_slide_v2(prs: Presentation, spec: SlideSpec) -> None:
    """Content slide with proper multi-item List handling."""
    layout = _safe_layout(prs, _LAYOUT_TITLE_CONTENT)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.get("heading", "")

    # Find body placeholder (index 1)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break

    if body_ph is None:
        return

    tf = body_ph.text_frame
    tf.word_wrap = True
    _populate_body(tf, spec.get("blocks", []))


def _add_table_slide(prs: Presentation, spec: SlideSpec) -> None:
    layout = _safe_layout(prs, _LAYOUT_TITLE_CONTENT)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.get("heading", "")

    table: Table = spec["table"]
    n_rows = len(table.rows) + 1  # +1 for header
    n_cols = len(table.headers) if table.headers else 1

    if n_rows < 1 or n_cols < 1:
        return

    # Position the table in the slide body area
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.0)
    height = Inches(0.5 * n_rows)

    pptx_table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Headers
    for col_idx, header_text in enumerate(table.headers):
        cell = pptx_table.cell(0, col_idx)
        cell.text = header_text
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True

    # Data rows
    for row_idx, row_data in enumerate(table.rows):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < n_cols:
                pptx_table.cell(row_idx + 1, col_idx).text = cell_text


def _add_figure_slide(prs: Presentation, spec: SlideSpec) -> None:
    layout = _safe_layout(prs, _LAYOUT_TITLE_CONTENT)
    slide = prs.slides.add_slide(layout)
    figure: Figure = spec["figure"]

    if slide.shapes.title:
        slide.shapes.title.text = spec.get("heading", "")

    src = figure.src
    if src and isfile(src):
        try:
            left = Inches(1.5)
            top = Inches(1.5)
            width = Inches(10.0)
            slide.shapes.add_picture(src, left, top, width=width)
        except Exception:
            logger.warning(
                "pptx_render: could not embed figure %s from %s",
                figure.id,
                src,
                exc_info=True,
            )

    if figure.caption:
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(10.0), Inches(0.8))
        tf = txBox.text_frame
        para = tf.paragraphs[0]
        para.text = figure.caption
        para.font.italic = True


def _add_code_slide(prs: Presentation, spec: SlideSpec) -> None:
    layout = _safe_layout(prs, _LAYOUT_TITLE_CONTENT)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.get("heading", "")

    code: Code = spec["code"]
    text = f"[{code.lang}]\n{code.text}" if code.lang else code.text

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.0)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = text
    para.font.name = "Courier New"
    para.font.size = Pt(14)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def render_pptx(doc: Document) -> bytes:
    """Render the agnostic Document tree to .pptx bytes via python-pptx.

    Opens a fresh :class:`~pptx.Presentation` (16:9 widescreen), collects a
    flat list of slide specs by walking the document tree, then emits slides
    in order.  If ``doc.meta.render_hints`` contains a ``max_slides`` key the
    slide list is truncated to that count.  The result is serialised to a
    :class:`~io.BytesIO` buffer and returned as raw bytes.
    """
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT

    specs = _collect_specs(doc)

    # Apply max_slides cap
    max_slides = doc.meta.render_hints.get("max_slides")
    if max_slides is not None:
        try:
            max_slides = int(max_slides)
            if len(specs) > max_slides:
                logger.warning(
                    "pptx_render: truncating from %d to %d slides (max_slides cap)",
                    len(specs),
                    max_slides,
                )
                specs = specs[:max_slides]
        except (TypeError, ValueError):
            logger.warning(
                "pptx_render: invalid max_slides value %r — ignoring", max_slides
            )

    # Emit slides
    for spec in specs:
        kind = spec["kind"]
        if kind == _TITLE_SLIDE:
            _add_title_slide(prs, spec)
        elif kind == _SECTION_DIVIDER:
            _add_section_divider(prs, spec)
        elif kind == _CONTENT_SLIDE:
            _add_content_slide_v2(prs, spec)
        elif kind == _TABLE_SLIDE:
            _add_table_slide(prs, spec)
        elif kind == _FIGURE_SLIDE:
            _add_figure_slide(prs, spec)
        elif kind == _CODE_SLIDE:
            _add_code_slide(prs, spec)
        else:
            logger.warning("pptx_render: unknown spec kind %s — skipped", kind)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
