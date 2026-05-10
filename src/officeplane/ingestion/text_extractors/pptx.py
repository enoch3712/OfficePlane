"""Extract per-slide text from a PPTX using python-pptx."""
from __future__ import annotations

import io

from pptx import Presentation


def extract_pptx_text(data: bytes) -> list[dict]:
    prs = Presentation(io.BytesIO(data))
    pages: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        title = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text or ""
            if shape.has_text_frame and shape == slide.shapes.title:
                title = text
            elif text.strip():
                chunks.append(text)
        body = "\n".join(chunks)
        full = f"# {title}\n\n{body}" if title else body
        pages.append({"page_number": i, "text": full.strip()})
    return pages
