"""Tests for the python-pptx renderer (Phase 10C)."""
from __future__ import annotations

import io
import json

from pptx import Presentation

from officeplane.content_agent.renderers.blocks import (
    Block,
    BlocksDocument,
    Lesson,
    Module,
)
from officeplane.content_agent.renderers.pptx_blocks import render_pptx


def test_renders_cover_section_and_content_slides():
    doc = BlocksDocument(
        title="Test Deck",
        modules=[
            Module(
                id="m1",
                title="Section One",
                lessons=[
                    Lesson(
                        id="l1",
                        title="Lesson Alpha",
                        blocks=[
                            Block(type="title", content="Slide A", order=0),
                            Block(type="text", content="bullet a1", order=1),
                            Block(type="text", content="bullet a2", order=2),
                            Block(type="title", content="Slide B", order=3),
                            Block(type="table", content=json.dumps({
                                "headers": ["x", "y"], "rows": [["1", "2"]],
                            }), order=4),
                        ],
                    )
                ],
            )
        ],
    )
    blob = render_pptx(doc)
    assert blob.startswith(b"PK")

    prs = Presentation(io.BytesIO(blob))
    titles = []
    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.text:
            titles.append(slide.shapes.title.text)
    # Cover (Test Deck), section (Section One), lesson title slide (Lesson Alpha),
    # 2 content slides (Slide A + Slide B)
    assert "Test Deck" in titles
    assert "Section One" in titles
    assert "Lesson Alpha" in titles
    assert "Slide A" in titles
    assert "Slide B" in titles

    # Slide B should have a table
    has_table = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                has_table = True
    assert has_table


def test_empty_blocks_document_renders_minimal_pptx():
    blob = render_pptx(BlocksDocument())
    assert blob.startswith(b"PK")
    prs = Presentation(io.BytesIO(blob))
    assert len(prs.slides) >= 1  # cover slide at minimum
