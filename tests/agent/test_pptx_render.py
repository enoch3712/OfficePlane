"""Tests for the PPTX tree renderer (Task 4)."""
import io

from pptx import Presentation

from officeplane.content_agent.renderers.document import parse_document
from officeplane.content_agent.renderers.pptx_render import render_pptx


def test_render_pptx_one_slide_per_lowest_section():
    doc = parse_document(
        {
            "type": "document",
            "meta": {"title": "BP Care"},
            "children": [
                {
                    "type": "section",
                    "level": 1,
                    "heading": "Why measure BP",
                    "children": [
                        {"type": "paragraph", "text": "Detect hypertension early."}
                    ],
                },
                {
                    "type": "section",
                    "level": 1,
                    "heading": "How to measure",
                    "children": [
                        {
                            "type": "section",
                            "level": 2,
                            "heading": "Cuff selection",
                            "children": [
                                {"type": "paragraph", "text": "Pick correct size."}
                            ],
                        },
                        {
                            "type": "section",
                            "level": 2,
                            "heading": "Positioning",
                            "children": [
                                {"type": "paragraph", "text": "Arm at heart level."}
                            ],
                        },
                    ],
                },
            ],
        }
    )
    blob = render_pptx(doc)
    assert isinstance(blob, bytes) and len(blob) > 5000
    pres = Presentation(io.BytesIO(blob))
    assert len(pres.slides) >= 4
    titles = [s.shapes.title.text for s in pres.slides if s.shapes.title]
    assert any("BP Care" in t for t in titles)
    assert any("Cuff selection" in t for t in titles)


def test_render_pptx_respects_max_slides_cap():
    children = [
        {
            "type": "section",
            "level": 1,
            "heading": f"S{i}",
            "children": [{"type": "paragraph", "text": "x"}],
        }
        for i in range(20)
    ]
    doc = parse_document(
        {
            "type": "document",
            "meta": {"title": "T", "render_hints": {"max_slides": 5}},
            "children": children,
        }
    )
    pres = Presentation(io.BytesIO(render_pptx(doc)))
    assert len(pres.slides) <= 5


def test_render_pptx_embeds_figure_from_prompt(tmp_path, monkeypatch):
    monkeypatch.setenv("OFFICEPLANE_IMAGE_PROVIDER", "placeholder")
    from officeplane.content_agent.renderers.document import parse_document
    from officeplane.content_agent.renderers.pptx_render import render_pptx

    doc = parse_document({"type": "document", "children": [
        {"type": "section", "level": 1, "heading": "BP", "children": [
            {"type": "figure", "id": "f2", "prompt": "labeled blood pressure cuff", "caption": "Fig 1"}
        ]}
    ]})
    blob = render_pptx(doc, workspace_dir=tmp_path)
    assert isinstance(blob, bytes) and len(blob) > 5000
    # Confirm the image file was generated
    assert (tmp_path / "images" / "f2.png").exists()


def test_render_pptx_emits_table_slide():
    doc = parse_document(
        {
            "type": "document",
            "meta": {"title": "T"},
            "children": [
                {
                    "type": "section",
                    "level": 1,
                    "heading": "Numbers",
                    "children": [
                        {
                            "type": "table",
                            "headers": ["A", "B"],
                            "rows": [["1", "2"], ["3", "4"]],
                        }
                    ],
                }
            ],
        }
    )
    pres = Presentation(io.BytesIO(render_pptx(doc)))
    # at least one slide has a native pptx table
    has_table = False
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_table:
                has_table = True
                break
    assert has_table
