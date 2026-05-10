"""Integration test for the generate-pptx-blocks handler (Phase 10C)."""
from __future__ import annotations

import io
import json
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation
from prisma import Prisma


@pytest.fixture
async def seeded_source_doc():
    db = Prisma()
    await db.connect()
    try:
        d = await db.document.create(data={
            "title": "Pptx Source",
            "summary": "Stub summary",
            "topics": ["t1"],
        })
        chap = await db.chapter.create(data={
            "documentId": d.id, "title": "C1", "orderIndex": 0,
        })
        await db.section.create(data={
            "chapterId": chap.id, "title": "S1", "orderIndex": 0,
        })
        yield d.id
        await db.document.delete(where={"id": d.id})
    finally:
        await db.disconnect()


@pytest.mark.asyncio
async def test_handler_produces_real_pptx_bytes(seeded_source_doc):
    from officeplane.content_agent.skill_executor import SkillExecutor

    ex = SkillExecutor()

    fake_blocks = {
        "schema_version": "1.0",
        "title": "Generated Deck",
        "modules": [{
            "id": "m1",
            "title": "Overview",
            "lessons": [{
                "id": "l1",
                "title": "Highlights",
                "order": 0,
                "blocks": [
                    {"type": "title", "content": "Topic 1", "order": 0},
                    {"type": "text", "content": "Key point one.", "order": 1},
                    {"type": "title", "content": "Topic 2", "order": 2},
                    {"type": "text", "content": "Key point two.", "order": 3},
                ],
            }],
        }],
    }
    fake_response = MagicMock()
    fake_response.choices = [MagicMock(message=MagicMock(content=json.dumps(fake_blocks)))]

    with patch("litellm.acompletion", new=AsyncMock(return_value=fake_response)):
        with patch(
            "officeplane.content_agent.skill_executor.SkillExecutor._emit_audit",
            new=AsyncMock(),
        ):
            output = await ex.invoke(
                "generate-pptx-blocks",
                {
                    "source_document_ids": [seeded_source_doc],
                    "brief": "intro deck",
                    "target_slide_count": 6,
                },
            )

    assert output["title"] == "Generated Deck"
    assert seeded_source_doc in output["source_document_ids"]

    from pathlib import Path
    blob = Path(output["file_path"]).read_bytes()
    assert blob.startswith(b"PK")
    prs = Presentation(io.BytesIO(blob))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title and s.shapes.title.text]
    assert "Generated Deck" in titles
    assert "Topic 1" in titles
    assert "Topic 2" in titles
