"""
Generate the OfficePlane pitch deck.
Run: python3 generate_deck.py
Output: OfficePlane_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ──────────────────────────────────────────────
BG_DARK    = RGBColor(0x06, 0x0A, 0x14)
BG_CARD    = RGBColor(0x0D, 0x12, 0x1F)
ORANGE     = RGBColor(0xE8, 0x62, 0x1C)
GREEN      = RGBColor(0x39, 0xFF, 0x14)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xC0, 0xC8, 0xD4)
MID_GRAY   = RGBColor(0x8A, 0x94, 0xA6)
DARK_GRAY  = RGBColor(0x3A, 0x42, 0x54)
RED_SOFT   = RGBColor(0xFF, 0x5F, 0x57)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # blank

# ── Helpers ───────────────────────────────────────────────────

def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    else:
        shape.adjustments[0] = 0.02
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_para(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), align=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p

def add_bullet(tf, text, size=16, color=LIGHT_GRAY, bullet_color=GREEN, indent=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_before = Pt(8)
    p.level = indent
    return p

def add_accent_line(slide, left, top, width, color=ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_green_dot(slide, left, top):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Pt(10), Pt(10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GREEN
    dot.line.fill.background()
    return dot

def add_number_badge(slide, left, top, number, bg_color=ORANGE):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = bg_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    return circle

# ── Image paths ───────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
PARENT = os.path.dirname(BASE)
IMG_PRETTIFY = os.path.join(BASE, "public", "diagram-prettify.png")
IMG_CONTRACT = os.path.join(BASE, "public", "diagram-contract.png")
IMG_DASHBOARD = os.path.join(PARENT, ".playwright-mcp", "dashboard-working.png")


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)

# Accent bar at top
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, color=ORANGE)

# Logo area
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(0.7), Inches(0.7), color=RGBColor(0x30, 0x18, 0x08))
add_text(slide, Inches(0.8), Inches(1.53), Inches(0.7), Inches(0.7), "{ }", size=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.7), Inches(1.5), Inches(4), Inches(0.5), "OfficePlane", size=28, color=WHITE, bold=True)
add_text(slide, Inches(1.7), Inches(1.95), Inches(4), Inches(0.4), "Agentic Runtime", size=14, color=MID_GRAY)

# Main title
add_text(slide, Inches(0.8), Inches(2.8), Inches(7.5), Inches(1.8),
         "The open-source agentic runtime\nfor deterministic document workflows",
         size=44, color=WHITE, bold=True)

# Subtitle
add_text(slide, Inches(0.8), Inches(4.8), Inches(8), Inches(0.8),
         "Give your AI agents a proper framework to read, write, and transform\noffice files -- with contracts, verification, and rollback built in.",
         size=20, color=LIGHT_GRAY)

# Tags
tag_y = Inches(5.9)
tags = ["Open Source", "DOCX", "XLSX", "PPTX", "PDF", "Plan-Execute-Verify"]
tag_x = Inches(0.8)
for tag in tags:
    w = Inches(len(tag) * 0.12 + 0.4)
    shape = add_shape_bg(slide, tag_x, tag_y, w, Inches(0.38), color=BG_CARD)
    shape.line.color.rgb = DARK_GRAY
    shape.line.width = Pt(1)
    add_text(slide, tag_x, tag_y, w, Inches(0.38), tag, size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    tag_x += w + Inches(0.15)

# Right side - run preview card
card_l, card_t = Inches(8.8), Inches(1.8)
card_w, card_h = Inches(4), Inches(4.8)
add_shape_bg(slide, card_l, card_t, card_w, card_h, color=BG_CARD)

add_text(slide, card_l + Inches(0.3), card_t + Inches(0.25), Inches(2.5), Inches(0.3),
         "Run #4,821  |  Completed", size=11, color=MID_GRAY)

# Green SUCCESS badge
badge = add_shape_bg(slide, card_l + Inches(2.7), card_t + Inches(0.2), Inches(1), Inches(0.32), color=RGBColor(0x0A, 0x2E, 0x05))
add_text(slide, card_l + Inches(2.7), card_t + Inches(0.2), Inches(1), Inches(0.32), "SUCCESS", size=10, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Stats row
stats = [("24", "Actions"), ("1.1s", "Latency"), ("42", "Pages")]
sx = card_l + Inches(0.3)
for val, label in stats:
    add_text(slide, sx, card_t + Inches(0.8), Inches(1), Inches(0.4), val, size=24, color=WHITE, bold=True)
    add_text(slide, sx, card_t + Inches(1.2), Inches(1), Inches(0.3), label, size=10, color=MID_GRAY)
    sx += Inches(1.15)

# Action steps
steps = [
    "Prettify -> extract structure         0.2s",
    "Contract Request -> plan              0.4s",
    "File assembly via LibreOffice         0.3s",
    "Verify output against contract        0.2s",
]
sy = card_t + Inches(1.8)
for step in steps:
    add_green_dot(slide, card_l + Inches(0.3), sy + Pt(4))
    add_text(slide, card_l + Inches(0.55), sy, Inches(3.2), Inches(0.28), step, size=10, color=MID_GRAY, font_name="Courier New")
    sy += Inches(0.32)

# Verification badge
vbox = add_shape_bg(slide, card_l + Inches(0.2), card_t + Inches(3.5), Inches(3.6), Inches(0.55), color=RGBColor(0x0A, 0x2E, 0x05))
add_text(slide, card_l + Inches(0.5), card_t + Inches(3.52), Inches(3.2), Inches(0.55),
         "Verification passed  |  Confidence: 0.97", size=12, color=GREEN, bold=True)

# GitHub url
add_text(slide, Inches(0.8), Inches(6.7), Inches(5), Inches(0.4),
         "github.com/enoch3712/AgenticDocs", size=14, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.7),
         "The Problem", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(1.2),
         "Agents are everywhere.\nBut they still can't reliably edit a spreadsheet.",
         size=38, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(0.6),
         "The agentic AI market is exploding -- yet there's no proper OSS framework for file manipulation.",
         size=18, color=LIGHT_GRAY)

# Two-column comparison
# Left: The Chaos
left_card = Inches(0.8)
card_top = Inches(3.5)
cw = Inches(5.5)
ch = Inches(3.2)

add_shape_bg(slide, left_card, card_top, cw, ch, color=BG_CARD)
add_text(slide, left_card + Inches(0.3), card_top + Inches(0.2), Inches(4), Inches(0.4),
         "Traditional MCP Approach", size=18, color=RED_SOFT, bold=True)

problems = [
    "Dozens of individual tool calls per document edit",
    "No atomicity -- partial failures leave corrupted files",
    "No verification -- agents can't confirm what they built",
    "No rollback -- mistakes are permanent",
    "Chaotic call graphs that are impossible to debug",
]
txBox = add_text(slide, left_card + Inches(0.3), card_top + Inches(0.7), Inches(5), Inches(2.3), "", size=14)
tf = txBox.text_frame
for i, prob in enumerate(problems):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {prob}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

# Right: What's needed
right_card = Inches(6.8)
add_shape_bg(slide, right_card, card_top, cw, ch, color=BG_CARD)
add_text(slide, right_card + Inches(0.3), card_top + Inches(0.2), Inches(4), Inches(0.4),
         "What Agents Actually Need", size=18, color=GREEN, bold=True)

needs = [
    "Declare intent, not individual steps",
    "Atomic execution -- all or nothing",
    "Built-in verification with confidence scores",
    "Transaction rollback on failure",
    "Deterministic, auditable action contracts",
]
txBox = add_text(slide, right_card + Inches(0.3), card_top + Inches(0.7), Inches(5), Inches(2.3), "", size=14)
tf = txBox.text_frame
for i, need in enumerate(needs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {need}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Market Opportunity", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(1.2),
         "Every AI agent needs to manipulate files.\nNo one has built the infrastructure layer -- until now.",
         size=38, color=WHITE, bold=True)

# Stats cards
stats_data = [
    ("82%", "of enterprise data lives in\noffice documents (Gartner)"),
    ("$28B+", "Agentic AI market by 2028\n(Markets & Markets)"),
    ("10x", "Growth in AI agent frameworks\nin the last 12 months"),
    ("0", "Open-source runtimes purpose-built\nfor agentic file manipulation"),
]

sx = Inches(0.8)
for val, desc in stats_data:
    cw_stat = Inches(2.8)
    add_shape_bg(slide, sx, Inches(3.2), cw_stat, Inches(2.2), color=BG_CARD)
    add_text(slide, sx + Inches(0.25), Inches(3.45), Inches(2.4), Inches(0.6),
             val, size=42, color=GREEN, bold=True)
    add_text(slide, sx + Inches(0.25), Inches(4.15), Inches(2.4), Inches(0.9),
             desc, size=13, color=LIGHT_GRAY)
    sx += cw_stat + Inches(0.2)

# Bottom insight
add_shape_bg(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.9), color=RGBColor(0x1A, 0x10, 0x05))
add_text(slide, Inches(1.1), Inches(5.85), Inches(11), Inches(0.9),
         "The gap: Agents can call APIs, search the web, write code -- but editing a Word doc or Excel sheet? That's still broken.\n"
         "OfficePlane fills this gap as the first open-source agentic runtime purpose-built for deterministic document workflows.",
         size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — INTRODUCING OFFICEPLANE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Introducing", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.9),
         "OfficePlane: the agentic runtime for deterministic document workflows",
         size=38, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(0.7),
         "An open-source framework that gives AI agents reliable, verified control over office documents.\n"
         "Plan. Contract. Apply. Verify. Every time.",
         size=18, color=LIGHT_GRAY)

# Three pillars
pillars = [
    ("Prettify", "Automatic structure\ndetection", "Point at any file. Get a clean,\nmachine-readable representation\nof chapters, sections, pages,\nand metadata.", GREEN),
    ("Contract Request", "Plan-build-verify\nin one pass", "Agent declares intent. OfficePlane\nplans, executes atomically, and\nverifies the output against the\ncontract.", ORANGE),
    ("Document Harness", "Orchestrate long-\nrunning workflows", "Context management, parallel\nexecution, checkpointing, and\nrollback. Plugs into any agent\nframework.", RGBColor(0x00, 0xBF, 0xFF)),
]

px = Inches(0.8)
for title, subtitle, desc, accent in pillars:
    pw = Inches(3.7)
    add_shape_bg(slide, px, Inches(3.5), pw, Inches(3.3), color=BG_CARD)
    # Accent line on card
    add_accent_line(slide, px + Inches(0.25), Inches(3.7), Inches(0.8), color=accent)
    add_text(slide, px + Inches(0.25), Inches(3.9), Inches(3.2), Inches(0.4),
             title, size=22, color=WHITE, bold=True)
    add_text(slide, px + Inches(0.25), Inches(4.3), Inches(3.2), Inches(0.5),
             subtitle, size=14, color=accent)
    add_text(slide, px + Inches(0.25), Inches(5.0), Inches(3.2), Inches(1.5),
             desc, size=13, color=LIGHT_GRAY)
    px += pw + Inches(0.3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS: CONTRACT MODEL (with image)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "How It Works", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.9),
         "The Contract Model: from chaos to determinism",
         size=38, color=WHITE, bold=True)

# Left column: explanation
add_text(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.5),
         "Traditional MCP: N chaotic tool calls", size=18, color=RED_SOFT, bold=True)
add_text(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(0.8),
         "edit_cell(), insert_row(), format_bold(), add_chart(),\n"
         "set_font(), merge_cells(), save()... hope it works.",
         size=14, color=MID_GRAY, font_name="Courier New")

add_text(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.5),
         "OfficePlane: 3 deterministic steps", size=18, color=GREEN, bold=True)

# Steps with numbers
steps_data = [
    ("1", "Plan", "Agent declares intent and schema"),
    ("2", "Contract", "Rules, validation, expected output locked in"),
    ("3", "Apply", "OfficePlane builds, verifies, delivers"),
]
sy = Inches(4.5)
for num, title, desc in steps_data:
    add_number_badge(slide, Inches(0.8), sy, num)
    add_text(slide, Inches(1.55), sy, Inches(2), Inches(0.35), title, size=18, color=WHITE, bold=True)
    add_text(slide, Inches(1.55), sy + Inches(0.32), Inches(4.5), Inches(0.3), desc, size=14, color=LIGHT_GRAY)
    sy += Inches(0.75)

# Right: Image
if os.path.exists(IMG_CONTRACT):
    slide.shapes.add_picture(IMG_CONTRACT, Inches(7), Inches(2.5), width=Inches(5.8))


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — PRETTIFY (with image)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Feature Deep-Dive", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.9),
         "Prettify: automatic document structure detection",
         size=38, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.3), Inches(6), Inches(0.8),
         "Point OfficePlane at any DOCX, PPTX, XLSX, or PDF.\n"
         "Prettify extracts the full document structure into\n"
         "a clean, machine-readable representation.",
         size=18, color=LIGHT_GRAY)

# Bullet features
features = [
    "Vision-powered extraction for complex layouts (Gemini API)",
    "Detects chapters, sections, pages, metadata automatically",
    "Works across all supported document formats",
    "Returns structured JSON your agent can reason about",
]
txBox = add_text(slide, Inches(0.8), Inches(3.6), Inches(5.5), Inches(2.5), "", size=14)
tf = txBox.text_frame
for i, feat in enumerate(features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {feat}"
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(10)

# Code snippet
code_box = add_shape_bg(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.3), color=BG_CARD)
add_text(slide, Inches(1.0), Inches(5.55), Inches(5.2), Inches(1.2),
         'doc = await plane.prettify("report.docx")\n'
         'print(doc.chapters)   # 8 chapters\n'
         'print(doc.sections)   # 23 sections\n'
         'print(doc.pages)      # 42 pages',
         size=13, color=GREEN, font_name="Courier New")

# Right: Image
if os.path.exists(IMG_PRETTIFY):
    slide.shapes.add_picture(IMG_PRETTIFY, Inches(7), Inches(2.5), width=Inches(5.8))


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — DOCUMENT HARNESS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Document Harness", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.9),
         "Like an agent harness, but for documents.\nFor your agentic flow.",
         size=38, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.5), Inches(6), Inches(0.7),
         "Read, prettify, plan, build, verify. All in one deterministic loop with\n"
         "checkpointing, parallel execution, and contract validation.",
         size=18, color=LIGHT_GRAY)

# Terminal mockup
term_l, term_t = Inches(0.8), Inches(3.5)
term_w, term_h = Inches(7.5), Inches(3.5)
add_shape_bg(slide, term_l, term_t, term_w, term_h, color=BG_CARD)

# Terminal dots
for i, c in enumerate([RED_SOFT, RGBColor(0xFE, 0xBC, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, term_l + Inches(0.2) + Inches(i * 0.22), term_t + Inches(0.15), Pt(10), Pt(10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()

add_text(slide, term_l + Inches(1), term_t + Inches(0.08), Inches(5), Inches(0.3),
         "OfficePlane Harness -- Q4 Board Deck Pipeline", size=10, color=MID_GRAY, font_name="Courier New")

terminal_lines = (
    "> /prettify-and-build is running...\n\n"
    "  Read(Q4-board-deck.docx)\n"
    "  -> Read 42 pages, 8 chapters detected\n\n"
    "  Prettify()\n"
    "  -> Extracted 8 chapters, 23 sections, 156 elements\n\n"
    "  ContractRequest(exec_summary)\n"
    "  -> Plan generated: 3 deliverables, 12 actions\n\n"
    "  Launching parallel document tasks...\n"
    "  -> Build executive summary (max 2 pages)\n"
    "  -> Extract KPI tables from chapters 3-5\n"
    "  -> Assemble final deck via LibreOffice driver"
)
add_text(slide, term_l + Inches(0.2), term_t + Inches(0.45), Inches(7), Inches(2.9),
         terminal_lines, size=11, color=LIGHT_GRAY, font_name="Courier New")

# Right side features
rx = Inches(9)
features = [
    ("Context Management", "Manages state across long-running\ndocument workflows"),
    ("Parallel Execution", "Multiple document tasks run\nconcurrently with progress tracking"),
    ("Atomic Rollback", "Never leave documents in a\nhalf-edited state"),
    ("Framework Agnostic", "Plugs into Claude, LangChain,\nCrewAI, or any custom setup"),
]
fy = Inches(3.5)
for title, desc in features:
    add_green_dot(slide, rx, fy + Pt(4))
    add_text(slide, rx + Inches(0.3), fy, Inches(3.5), Inches(0.3),
             title, size=16, color=WHITE, bold=True)
    add_text(slide, rx + Inches(0.3), fy + Inches(0.3), Inches(3.5), Inches(0.5),
             desc, size=12, color=MID_GRAY)
    fy += Inches(0.85)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — SHEETTOOLS: EXCEL FOR AGENTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "SheetTools", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.9),
         "Excel manipulation for agents.\nPlan-Execute-Rollback built in.",
         size=38, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.5),
         "Dependency-free XLSX editing with transactions, structured results, and full rollback support.",
         size=18, color=LIGHT_GRAY)

# Code example (left)
code_box = add_shape_bg(slide, Inches(0.8), Inches(3.3), Inches(6.2), Inches(3.8), color=BG_CARD)

# Terminal header
for i, c in enumerate([RED_SOFT, RGBColor(0xFE, 0xBC, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0) + Inches(i * 0.22), Inches(3.45), Pt(10), Pt(10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()

add_text(slide, Inches(1.8), Inches(3.35), Inches(3), Inches(0.3),
         "sheettools_example.py", size=10, color=MID_GRAY, font_name="Courier New")

code = (
    'from officeplane.sheettools import (\n'
    '    SpreadsheetEditor, SpreadsheetPlan, PlanExecutor\n'
    ')\n\n'
    'with SpreadsheetEditor("sales.xlsx") as editor:\n'
    '    plan = SpreadsheetPlan("Q4 Sales Report")\n'
    '    plan.add_sheet("Data")\n'
    '    plan.create_table(\n'
    '        "Data", "A1",\n'
    '        headers=["Month", "Revenue", "Growth"],\n'
    '        rows=[["Oct", 42000, "12%"],\n'
    '              ["Nov", 48000, "14%"],\n'
    '              ["Dec", 55000, "15%"]]\n'
    '    )\n\n'
    '    executor = PlanExecutor(editor)\n'
    '    result = executor.execute(plan)\n'
    '    # Atomic! Rolls back on any failure.'
)
add_text(slide, Inches(1.0), Inches(3.7), Inches(5.8), Inches(3.2),
         code, size=12, color=GREEN, font_name="Courier New")

# Right side: capabilities
rx = Inches(7.5)
caps = [
    ("SpreadsheetEditor", "Open, edit, save XLSX files.\nPure Python -- no openpyxl dependency.\nDirect OOXML manipulation."),
    ("Transaction Support", "with editor.transaction() as tx:\nAutomatic rollback on exception.\nDeep-copy backup & restore."),
    ("SheetReader", "list_sheets(), get_row(), get_column()\nget_used_range(), get_sheet_summary()"),
    ("CellModifier", "set_values(), set_formula(),\nclear_values(), replace_in_sheet()"),
    ("TableBuilder", "create_table() with headers & rows.\nadd_totals_row() with SUM formulas."),
    ("PlanExecutor", "Validate -> Execute -> Verify.\nDry-run mode. Step callbacks.\nFull result tracking per step."),
]

fy = Inches(1.2)
for title, desc in caps:
    add_green_dot(slide, rx, fy + Pt(4))
    add_text(slide, rx + Inches(0.3), fy, Inches(5), Inches(0.3),
             title, size=14, color=WHITE, bold=True)
    add_text(slide, rx + Inches(0.3), fy + Inches(0.28), Inches(5), Inches(0.65),
             desc, size=11, color=MID_GRAY, font_name="Courier New")
    fy += Inches(0.95)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Architecture", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.5),
         "Modular, extensible, production-ready",
         size=38, color=WHITE, bold=True)

# Architecture blocks
modules = [
    ("DocTools", "Word/DOCX editing\nwith transactions", ORANGE, Inches(0.8), Inches(2.3)),
    ("SheetTools", "Excel/XLSX editing\nwith plan-execute", GREEN, Inches(3.5), Inches(2.3)),
    ("Ingestion", "Vision-powered\nstructure extraction", RGBColor(0x00, 0xBF, 0xFF), Inches(6.2), Inches(2.3)),
    ("Drivers", "LibreOffice | Rust\nMock | Custom", RGBColor(0xFF, 0xD7, 0x00), Inches(8.9), Inches(2.3)),

    ("Planning", "LLM-based action\ngraph generation", ORANGE, Inches(0.8), Inches(4.1)),
    ("Management", "Instances, tasks\nqueue, orchestration", GREEN, Inches(3.5), Inches(4.1)),
    ("Memory/RAG", "Vector store\nembeddings, search", RGBColor(0x00, 0xBF, 0xFF), Inches(6.2), Inches(4.1)),
    ("Observability", "Logging, metrics\nPrometheus, traces", RGBColor(0xFF, 0xD7, 0x00), Inches(8.9), Inches(4.1)),
]

mw, mh = Inches(2.4), Inches(1.5)
for name, desc, accent, mx, my in modules:
    add_shape_bg(slide, mx, my, mw, mh, color=BG_CARD)
    add_accent_line(slide, mx + Inches(0.15), my + Inches(0.12), Inches(0.6), color=accent)
    add_text(slide, mx + Inches(0.15), my + Inches(0.28), Inches(2.1), Inches(0.35),
             name, size=16, color=WHITE, bold=True)
    add_text(slide, mx + Inches(0.15), my + Inches(0.65), Inches(2.1), Inches(0.7),
             desc, size=12, color=MID_GRAY)

# Bottom layer: API + Storage + UI
bottom_modules = [
    ("FastAPI", "REST + WebSocket\nReal-time updates"),
    ("PostgreSQL + pgvector", "Documents, vectors\nPrisma ORM"),
    ("Next.js Dashboard", "Management UI\nTask monitoring"),
    ("Docker", "One-command deploy\nAll-in-one container"),
]

bx = Inches(0.8)
bw = Inches(2.85)
for name, desc in bottom_modules:
    add_shape_bg(slide, bx, Inches(5.9), bw, Inches(1.1), color=RGBColor(0x12, 0x18, 0x28))
    add_text(slide, bx + Inches(0.15), Inches(5.95), Inches(2.6), Inches(0.3),
             name, size=14, color=ORANGE, bold=True)
    add_text(slide, bx + Inches(0.15), Inches(6.25), Inches(2.6), Inches(0.55),
             desc, size=11, color=MID_GRAY)
    bx += bw + Inches(0.15)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — CODE EXAMPLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Developer Experience", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.5),
         "One intent. Full lifecycle. Zero guesswork.",
         size=38, color=WHITE, bold=True)

# Large code block
code_box = add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(7), Inches(4.8), color=BG_CARD)
for i, c in enumerate([RED_SOFT, RGBColor(0xFE, 0xBC, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0) + Inches(i * 0.22), Inches(2.35), Pt(10), Pt(10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()
add_text(slide, Inches(1.8), Inches(2.25), Inches(3), Inches(0.3),
         "example.py", size=10, color=MID_GRAY, font_name="Courier New")

code = (
    'from officeplane import OfficePlane\n\n'
    '# Initialize with a driver\n'
    'plane = OfficePlane(driver="libreoffice")\n\n'
    '# Prettify: extract document structure\n'
    'doc = await plane.prettify("report.docx")\n'
    'print(doc.chapters, doc.sections)\n\n'
    '# Contract Request: agent defines the plan\n'
    'contract = plane.contract_request(\n'
    '    schema="exec_summary",\n'
    '    source=doc,\n'
    '    rules=["max 2 pages", "include KPIs"]\n'
    ')\n\n'
    '# OfficePlane builds the file\n'
    'result = await plane.execute(contract)\n'
    'print(f"Confidence: {result.confidence}")\n'
    '# Confidence: 0.97'
)
add_text(slide, Inches(1.0), Inches(2.65), Inches(6.5), Inches(4.2),
         code, size=14, color=GREEN, font_name="Courier New")

# Right: key principles
rx = Inches(8.5)
principles = [
    ("Declare intent,\nnot steps", "Let the runtime figure out\nthe action graph. You say\nwhat, OfficePlane says how."),
    ("Atomic execution\nwith rollback", "Never leave documents in a\nhalf-edited state. Every\noperation is all-or-nothing."),
    ("Built-in\nverification", "Confidence scoring and\noutcome validation included.\nKnow what you got."),
    ("Result types,\nnot exceptions", "Ok/Err pattern throughout.\nStructured error codes.\nNo surprises."),
]

fy = Inches(2.3)
for title, desc in principles:
    add_text(slide, rx, fy, Inches(4), Inches(0.5), title, size=16, color=WHITE, bold=True)
    add_text(slide, rx, fy + Inches(0.48), Inches(4), Inches(0.7), desc, size=12, color=MID_GRAY)
    fy += Inches(1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — OPEN SOURCE + BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Open-Source Strategy", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.9),
         "Open-source core. Commercial execution.",
         size=38, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(0.5),
         "Self-host free forever. The moat is deterministic edit contracts + fidelity verification, not file storage.",
         size=18, color=LIGHT_GRAY)

# Two columns
# Left: OSS Core
add_shape_bg(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(3.8), color=BG_CARD)
add_text(slide, Inches(1.1), Inches(3.35), Inches(5), Inches(0.4),
         "Open Source (Permissive License)", size=18, color=GREEN, bold=True)

oss_items = [
    "Edit Program IR spec + validators + compiler",
    "Deterministic execution engine (runs, events, retries)",
    "Driver interfaces + reference drivers (DOCX, XLSX)",
    "Verification toolkit (render + visual diff + structural checks)",
    "SDKs + CLI + local docker-compose runner",
    "SheetTools, DocTools, Ingestion pipeline",
]
txBox = add_text(slide, Inches(1.1), Inches(3.85), Inches(5), Inches(3), "", size=14)
tf = txBox.text_frame
for i, item in enumerate(oss_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

# Right: Commercial
add_shape_bg(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.8), color=BG_CARD)
add_text(slide, Inches(7.1), Inches(3.35), Inches(5), Inches(0.4),
         "Commercial (Managed + Enterprise)", size=18, color=ORANGE, bold=True)

commercial_items = [
    "Hosted runtime: scalable render farm, queues, observability",
    "Enterprise: SSO/SAML, SCIM, RBAC, approval workflows",
    "Policy packs, private networking, SLAs",
    "Commercial drivers: high-fidelity rendering engines",
    "Specialized verifiers, regulated connectors",
    "Volume discounts, dedicated infrastructure",
]
txBox = add_text(slide, Inches(7.1), Inches(3.85), Inches(5), Inches(3), "", size=14)
tf = txBox.text_frame
for i, item in enumerate(commercial_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — ENTERPRISE INTEGRATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Enterprise Integration", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.9),
         "Works alongside your existing ECM.\nNo rip-and-replace required.",
         size=38, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.5), Inches(10), Inches(0.5),
         "OfficePlane runs as a sidecar to Alfresco, SharePoint, or any document management system.\n"
         "Customers keep their governance. You add deterministic agent capabilities.",
         size=18, color=LIGHT_GRAY)

# Two-column table
# Left: ECM (System of Record)
add_shape_bg(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.2), color=BG_CARD)
add_text(slide, Inches(1.1), Inches(3.65), Inches(5), Inches(0.4),
         "ECM  --  System of Record", size=18, color=RGBColor(0x00, 0xBF, 0xFF), bold=True)

ecm_items = [
    "Stores canonical documents and versions",
    "Permissions / RBAC and access control",
    "Metadata, folders, lifecycle management",
    "Audit logs (repository events)",
    "Integration surface (REST / CMIS)",
]
txBox = add_text(slide, Inches(1.1), Inches(4.2), Inches(5), Inches(2.3), "", size=14)
tf = txBox.text_frame
for i, item in enumerate(ecm_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)

# Right: OfficePlane (System of Change)
add_shape_bg(slide, Inches(6.8), Inches(3.5), Inches(5.5), Inches(3.2), color=BG_CARD)
add_text(slide, Inches(7.1), Inches(3.65), Inches(5), Inches(0.4),
         "OfficePlane  --  System of Change + Proof", size=18, color=GREEN, bold=True)

op_items = [
    "Executes Plan -> Contract -> Apply -> Verify runs",
    "Policy-aware contracts (what may be edited)",
    "IR compiler + deterministic drivers (DOCX/XLSX/PDF)",
    "Run ledger: immutable event history + provenance",
    "Verification bundle: visual diff + structural checks",
]
txBox = add_text(slide, Inches(7.1), Inches(4.2), Inches(5), Inches(2.3), "", size=14)
tf = txBox.text_frame
for i, item in enumerate(op_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(8)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — PRICING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.5))

add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
         "Pricing", size=14, color=ORANGE, bold=True)
add_text(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.5),
         "Self-host free forever. Managed service starts at $0.",
         size=38, color=WHITE, bold=True)

# Pricing cards
tiers = [
    ("Free", "$0", "", [
        "1,000 Actions",
        "500 Prettify pages",
        "50 Contract Requests",
        "1 GB Storage",
        "LibreOffice driver",
        "Community support",
    ], DARK_GRAY),
    ("Developer", "$9", "/month", [
        "25,000 Actions",
        "10,000 Prettify pages",
        "1,000 Contract Requests",
        "10 GB Storage",
        "All drivers",
        "Email support",
    ], GREEN),
    ("Scaler", "$49", "/month", [
        "200,000 Actions",
        "100,000 Prettify pages",
        "10,000 Contract Requests",
        "50 GB Storage",
        "All drivers",
        "Priority support",
    ], GREEN),
    ("Pro", "$199", "/month", [
        "1,000,000 Actions",
        "500,000 Prettify pages",
        "100,000 Contract Requests",
        "200 GB Storage",
        "All + custom drivers",
        "Dedicated support",
    ], GREEN),
]

tx = Inches(0.8)
tier_w = Inches(2.85)
for name, price, period, features_list, accent in tiers:
    add_shape_bg(slide, tx, Inches(2.3), tier_w, Inches(4.5), color=BG_CARD)
    add_text(slide, tx + Inches(0.2), Inches(2.45), Inches(2.4), Inches(0.3),
             name, size=14, color=accent, bold=True)
    price_text = price + (f" {period}" if period else "")
    add_text(slide, tx + Inches(0.2), Inches(2.75), Inches(2.4), Inches(0.5),
             price_text, size=28, color=WHITE, bold=True)

    fy = Inches(3.4)
    for feat in features_list:
        add_text(slide, tx + Inches(0.2), fy, Inches(2.5), Inches(0.28),
                 f"  {feat}", size=12, color=LIGHT_GRAY)
        fy += Inches(0.3)

    tx += tier_w + Inches(0.2)

# Enterprise bar
add_shape_bg(slide, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.45), color=RGBColor(0x12, 0x18, 0x28))
add_text(slide, Inches(1.1), Inches(6.95), Inches(5), Inches(0.45),
         "Enterprise: Volume discounts, SSO, dedicated infra, custom drivers, SLA",
         size=13, color=MID_GRAY, bold=True)
add_text(slide, Inches(9.5), Inches(6.95), Inches(3), Inches(0.45),
         "Talk to us ->", size=13, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — CALL TO ACTION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
dark_bg(slide)

# Big orange accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ORANGE
bar.line.fill.background()

# Center content
add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
         "Ready to build the future\nof document automation?",
         size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(3.0), Inches(9), Inches(0.8),
         "Join the community building the agentic runtime for office workflows.\n"
         "Star the repo, contribute drivers, and shape the future of agentic document workflows.",
         size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# CTA buttons (as shapes)
btn1 = add_shape_bg(slide, Inches(3.5), Inches(4.3), Inches(3), Inches(0.7), color=ORANGE)
add_text(slide, Inches(3.5), Inches(4.32), Inches(3), Inches(0.7),
         "Star on GitHub", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

btn2 = add_shape_bg(slide, Inches(7), Inches(4.3), Inches(3), Inches(0.7), color=BG_CARD)
btn2.line.color.rgb = DARK_GRAY
btn2.line.width = Pt(1.5)
add_text(slide, Inches(7), Inches(4.32), Inches(3), Inches(0.7),
         "Read Documentation", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# URL
add_text(slide, Inches(2), Inches(5.4), Inches(9), Inches(0.5),
         "github.com/enoch3712/AgenticDocs", size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Bottom card: waitlist
add_shape_bg(slide, Inches(3.5), Inches(6.0), Inches(6.3), Inches(1), color=BG_CARD)
add_text(slide, Inches(3.8), Inches(6.05), Inches(5.8), Inches(0.35),
         "Get early access to the Scale tier", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(3.8), Inches(6.4), Inches(5.8), Inches(0.5),
         "Join the waitlist for hosted agentic runs with audit logs, team features, and premium support.\n"
         "Early access includes $200 in free agentic actions.",
         size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(BASE, "OfficePlane_Deck.pptx")
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
